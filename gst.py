import os
import numpy as np
from sentence_transformers import SentenceTransformer, util
from datasets import Dataset
import streamlit as st
from docx import Document
import PyPDF2
from pptx import Presentation
import torch
from transformers import pipeline 

def extract_text(file):
    ext = file.name.split('.')[-1].lower()

    if ext == 'docx':
        return extract_text_from_docx(file)
    elif ext == 'pdf':
        return extract_text_from_pdf(file)
    elif ext == 'pptx':
        return extract_text_from_pptx(file)
    else:
        st.error("Unsupported file format. Please upload a DOCX, PDF, or PPTX file.")
        return ""



# Extract text from DOCX file
def extract_text_from_docx(docx_file):
    try:
        doc = Document(docx_file)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text if text.strip() else "No readable text found in the DOCX file."
    except Exception as e:
        st.error(f"Error reading DOCX file: {str(e)}")
        return ""

# Extract text from PDF file
def extract_text_from_pdf(pdf_file):
    try:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page_num in range(len(pdf_reader.pages)):
            text += pdf_reader.pages[page_num].extract_text() + "\n"
        return text if text.strip() else "No readable text found in the PDF file."
    except Exception as e:
        st.error(f"Error reading PDF file: {str(e)}")
        return ""

# Extract text from PPTX file
def extract_text_from_pptx(pptx_file):
    try:
        presentation = Presentation(pptx_file)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text if text.strip() else "No readable text found in the PPTX file."
    except Exception as e:
        st.error(f"Error reading PPTX file: {str(e)}")
        return ""
    
# Split document into smaller chunks
def split_text_into_chunks(text, chunk_size=512):
    if not text.strip():
        return []
    words = text.split()
    chunks = [" ".join(words[i:i + chunk_size]) for i in range(0, len(words), chunk_size)]
    return [chunk for chunk in chunks if chunk.strip()]

# Retrieve the most relevant passage using cosine similarity
def retrieve_relevant_text(question, document_chunks, embedding_model):
    question_embedding = embedding_model.encode(question, convert_to_tensor=True)

    if not document_chunks:
        return [], []

    document_embeddings = embedding_model.encode(document_chunks, convert_to_tensor=True)
    similarities = util.pytorch_cos_sim(question_embedding, document_embeddings)[0]

    top_k = min(3, len(document_chunks))  # Get top 3 relevant chunks
    top_results = torch.topk(similarities, k=top_k)

    return [document_chunks[i] for i in top_results.indices], top_results.values




# Helper function to split text into chunks suitable for the model's max token length
def split_text_for_model(input_text, max_tokens=512):
    words = input_text.split()
    chunks = [" ".join(words[i:i + max_tokens]) for i in range(0, len(words), max_tokens)]
    return chunks

# Generate concise answer by summarizing the context in chunks
def generate_summary(context, question, max_tokens=512):
    summarizer = pipeline("text2text-generation", model="google/flan-t5-small")

    # Split the context into smaller chunks if it's too long for the model
    context_chunks = split_text_for_model(context, max_tokens)

    # Generate and combine summaries for each chunk
    summaries = []
    for chunk in context_chunks:
        input_text = f"Context: {chunk}\nQuestion: {question}\nSummarize the answer in a concise way."
        summary = summarizer(input_text, max_length=100, min_length=30, do_sample=False)[0]['generated_text']
        summaries.append(summary)

    # Combine the summaries
    return " ".join(summaries)

# Streamlit App - Main logic (Unchanged)
def main():
    st.set_page_config(page_title="GST Smart Guide", page_icon=":books:", layout="centered")
    st.title("📚 GST Smart Guide")

    # UI for document upload and question input
    file = st.file_uploader("Upload GST Guide (DOCX, PDF, PPTX)", type=["docx", "pdf", "pptx"])
    user_question = st.text_input("🔍 Enter your question:")
    submit = st.button("Submit")

    if file is not None:
        document_text = extract_text(file)
        document_chunks = split_text_into_chunks(document_text)

        if not document_chunks:
            st.error("The uploaded document contains no readable text.")
            return

        # Load Sentence Transformer model
        embedding_model = SentenceTransformer('sentence-transformers/all-mpnet-base-v2')

        if submit and user_question:
            with st.spinner("⏳ Retrieving the answer..."):
                relevant_texts, scores = retrieve_relevant_text(user_question, document_chunks, embedding_model)

                if relevant_texts:
                    st.success("✅ Answer Retrieved!")
                    
                    # Take the most relevant context and generate a concise answer
                    context = " ".join(relevant_texts)
                    generated_answer = generate_summary(context, user_question)
                    
                    # Display the final generated answer
                    st.markdown(f"### Answer:")
                    st.write(generated_answer)
                else:
                    st.error("No relevant texts found. Please try a different query or upload a different document.")

if __name__ == '__main__':
    main()

